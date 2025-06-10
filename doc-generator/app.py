# doc-generator/app.py

from flask import Flask, request, jsonify, send_file
import os
import json
import shutil
import subprocess
from datetime import datetime
from pathlib import Path
import redis
import requests
from typing import Dict, List, Any, Optional
import markdown
from jinja2 import Environment, FileSystemLoader
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import matplotlib.pyplot as plt
import seaborn as sns
from reportlab.lib.pagesizes import letter, A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, PageBreak, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib import colors

app = Flask(__name__)

# Redis connection
redis_client = redis.Redis(
    host=os.environ.get('REDIS_HOST', 'redis'),
    port=int(os.environ.get('REDIS_PORT', 6379)),
    decode_responses=True
)

# Template paths
TEMPLATES_DIR = Path('/templates')
OUTPUT_DIR = Path('/output')
WORKSPACE_DIR = Path('/workspace')

# Ensure directories exist
TEMPLATES_DIR.mkdir(exist_ok=True)
OUTPUT_DIR.mkdir(exist_ok=True)

class DocumentGenerator:
    def __init__(self):
        self.jinja_env = Environment(loader=FileSystemLoader(str(TEMPLATES_DIR)))
        self.setup_latex_templates()
        
    def setup_latex_templates(self):
        """Setup default LaTeX templates"""
        ieee_template = r"""
\documentclass[conference]{IEEEtran}
\IEEEoverridecommandlockouts
\usepackage{cite}
\usepackage{amsmath,amssymb,amsfonts}
\usepackage{algorithmic}
\usepackage{graphicx}
\usepackage{textcomp}
\usepackage{xcolor}
\usepackage{hyperref}
\usepackage{listings}
\usepackage{booktabs}

\lstset{
    basicstyle=\footnotesize\ttfamily,
    breaklines=true,
    numbers=left,
    numberstyle=\tiny,
    frame=single,
    tabsize=2,
    showstringspaces=false
}

\begin{document}

\title{{{ title }}}

\author{
    \IEEEauthorblockN{{{ author_name }}}
    \IEEEauthorblockA{
        {{ affiliation }}\\
        {{ email }}
    }
}

\maketitle

\begin{abstract}
{{ abstract }}
\end{abstract}

\begin{IEEEkeywords}
{{ keywords }}
\end{IEEEkeywords}

{{ content }}

\bibliographystyle{IEEEtran}
\bibliography{references}

\end{document}
"""
        
        # Save IEEE template
        ieee_path = TEMPLATES_DIR / 'ieee_template.tex'
        with open(ieee_path, 'w') as f:
            f.write(ieee_template)
    
    async def generate_pdf_report(self, session_id: str, research_data: Dict) -> str:
        """Generate PDF report from research results"""
        try:
            # Create output filename
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            filename = f"research_report_{session_id}_{timestamp}.pdf"
            filepath = OUTPUT_DIR / filename
            
            # Create PDF document
            doc = SimpleDocTemplate(
                str(filepath),
                pagesize=letter,
                rightMargin=72,
                leftMargin=72,
                topMargin=72,
                bottomMargin=18,
            )
            
            # Container for the 'Flowable' objects
            elements = []
            
            # Define styles
            styles = getSampleStyleSheet()
            title_style = ParagraphStyle(
                'CustomTitle',
                parent=styles['Title'],
                fontSize=24,
                textColor=colors.HexColor('#1a1a1a'),
                spaceAfter=30,
            )
            heading_style = ParagraphStyle(
                'CustomHeading',
                parent=styles['Heading1'],
                fontSize=16,
                textColor=colors.HexColor('#2c3e50'),
                spaceAfter=12,
            )
            
            # Title
            elements.append(Paragraph(research_data['topic'], title_style))
            elements.append(Spacer(1, 12))
            
            # Executive Summary
            elements.append(Paragraph("Executive Summary", heading_style))
            elements.append(Paragraph(research_data.get('synthesis', ''), styles['Normal']))
            elements.append(Spacer(1, 12))
            
            # Research Objectives
            elements.append(Paragraph("Research Objectives", heading_style))
            for obj in research_data.get('objectives', []):
                elements.append(Paragraph(f"• {obj}", styles['Normal']))
            elements.append(Spacer(1, 12))
            
            # Key Findings
            elements.append(Paragraph("Key Findings", heading_style))
            elements.append(Paragraph(research_data.get('synthesis', ''), styles['Normal']))
            elements.append(Spacer(1, 12))
            
            # Research Gaps
            if 'gaps' in research_data:
                elements.append(Paragraph("Identified Research Gaps", heading_style))
                for gap in research_data['gaps']:
                    elements.append(Paragraph(f"• {gap}", styles['Normal']))
                elements.append(Spacer(1, 12))
            
            # Proposed Approach
            if 'proposedApproach' in research_data:
                elements.append(Paragraph("Proposed Approach", heading_style))
                elements.append(Paragraph(research_data['proposedApproach'], styles['Normal']))
                elements.append(Spacer(1, 12))
            
            # Technical Details
            if 'technicalDetails' in research_data:
                elements.append(PageBreak())
                elements.append(Paragraph("Technical Implementation Details", heading_style))
                
                for detail in research_data['technicalDetails']:
                    elements.append(Paragraph(detail['aspect'], styles['Heading2']))
                    elements.append(Paragraph(detail['description'], styles['Normal']))
                    elements.append(Spacer(1, 6))
                    
                    if detail.get('implementation'):
                        elements.append(Paragraph("Implementation:", styles['Heading3']))
                        elements.append(Paragraph(detail['implementation'], styles['Code']))
                    elements.append(Spacer(1, 12))
            
            # Development Progress
            if 'development' in research_data:
                elements.append(Paragraph("Development Progress", heading_style))
                dev_data = research_data['development']
                
                # Create progress table
                progress_data = [
                    ['Component', 'Status', 'Completion'],
                    ['Code Implementation', dev_data.get('codeStatus', 'In Progress'), 
                     f"{dev_data.get('codeCompletion', 0)}%"],
                    ['Unit Tests', dev_data.get('testStatus', 'Pending'), 
                     f"{dev_data.get('testCoverage', 0)}%"],
                    ['Documentation', dev_data.get('docStatus', 'Pending'), 
                     f"{dev_data.get('docCompletion', 0)}%"],
                ]
                
                progress_table = Table(progress_data)
                progress_table.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0, 0), (-1, 0), 14),
                    ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                    ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                    ('GRID', (0, 0), (-1, -1), 1, colors.black)
                ]))
                
                elements.append(progress_table)
                elements.append(Spacer(1, 12))
            
            # Testing Results
            if 'testResults' in research_data:
                elements.append(Paragraph("Testing Results", heading_style))
                test_data = research_data['testResults']
                
                test_summary = f"""
                Total Tests: {test_data.get('total', 0)}
                Passed: {test_data.get('passed', 0)}
                Failed: {test_data.get('failed', 0)}
                Coverage: {test_data.get('coverage', 0)}%
                """
                
                elements.append(Paragraph(test_summary, styles['Normal']))
                elements.append(Spacer(1, 12))
            
            # References
            if 'papers' in research_data:
                elements.append(PageBreak())
                elements.append(Paragraph("References", heading_style))
                
                for i, paper in enumerate(research_data['papers'][:20], 1):
                    ref_text = f"[{i}] {', '.join(paper['authors'][:3])}. \"{paper['title']}\". "
                    if paper.get('venue'):
                        ref_text += f"{paper['venue']}, "
                    ref_text += f"{paper.get('year', 'n.d.')}."
                    
                    elements.append(Paragraph(ref_text, styles['Normal']))
                    elements.append(Spacer(1, 6))
            
            # Build PDF
            doc.build(elements)
            
            return str(filepath)
            
        except Exception as e:
            raise Exception(f"Error generating PDF report: {str(e)}")
    
    async def generate_latex_paper(
        self, 
        session_id: str, 
        research_data: Dict,
        template: str = 'ieee'
    ) -> str:
        """Generate LaTeX paper from research results"""
        try:
            # Create output directory
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            paper_dir = OUTPUT_DIR / f"paper_{session_id}_{timestamp}"
            paper_dir.mkdir(exist_ok=True)
            
            # Get template
            template_file = TEMPLATES_DIR / f"{template}_template.tex"
            if not template_file.exists():
                template_file = TEMPLATES_DIR / "ieee_template.tex"
            
            # Prepare paper content
            paper_data = {
                'title': research_data.get('topic', 'Untitled Research'),
                'author_name': 'Research Team',
                'affiliation': 'Research Institution',
                'email': 'research@institution.edu',
                'abstract': self._generate_abstract(research_data),
                'keywords': ', '.join(self._extract_keywords(research_data)),
                'content': self._generate_paper_content(research_data)
            }
            
            # Generate main tex file
            with open(template_file, 'r') as f:
                template_content = f.read()
            
            # Use jinja2 for template rendering
            from jinja2 import Template
            template = Template(template_content)
            latex_content = template.render(**paper_data)
            
            # Write main tex file
            main_tex = paper_dir / 'main.tex'
            with open(main_tex, 'w') as f:
                f.write(latex_content)
            
            # Generate bibliography
            self._generate_bibliography(paper_dir, research_data.get('papers', []))
            
            # Copy any figures
            figures_dir = paper_dir / 'figures'
            figures_dir.mkdir(exist_ok=True)
            
            # Compile LaTeX (attempt multiple times for references)
            for _ in range(3):
                subprocess.run(
                    ['pdflatex', '-interaction=nonstopmode', 'main.tex'],
                    cwd=paper_dir,
                    capture_output=True
                )
            
            # Run bibtex
            subprocess.run(['bibtex', 'main'], cwd=paper_dir, capture_output=True)
            
            # Final compilation
            subprocess.run(
                ['pdflatex', '-interaction=nonstopmode', 'main.tex'],
                cwd=paper_dir,
                capture_output=True
            )
            
            # Create zip archive
            archive_path = OUTPUT_DIR / f"paper_{session_id}_{timestamp}.zip"
            shutil.make_archive(
                str(archive_path.with_suffix('')),
                'zip',
                paper_dir
            )
            
            return str(archive_path)
            
        except Exception as e:
            raise Exception(f"Error generating LaTeX paper: {str(e)}")
    
    def _generate_abstract(self, research_data: Dict) -> str:
        """Generate abstract from research data"""
        # Call Ollama to generate abstract
        ollama_url = f"http://{os.environ.get('OLLAMA_HOST', 'ollama')}:11434/api/chat"
        
        prompt = f"""
        Generate a concise academic abstract (150-250 words) for a research paper based on:
        
        Topic: {research_data.get('topic')}
        Objectives: {research_data.get('objectives', [])}
        Key Findings: {research_data.get('synthesis', '')}
        Proposed Approach: {research_data.get('proposedApproach', '')}
        
        The abstract should follow the standard structure:
        1. Context/Background
        2. Problem statement
        3. Methodology
        4. Key findings
        5. Significance
        """
        
        try:
            response = requests.post(ollama_url, json={
                'model': os.environ.get('OLLAMA_MODEL', 'mixtral:8x7b'),
                'messages': [
                    {'role': 'system', 'content': 'You are writing an academic paper abstract.'},
                    {'role': 'user', 'content': prompt}
                ],
                'stream': False
            })
            
            return response.json()['message']['content']
        except:
            return "Abstract generation in progress..."
    
    def _extract_keywords(self, research_data: Dict) -> List[str]:
        """Extract keywords from research data"""
        keywords = set()
        
        # Extract from topic
        topic_words = research_data.get('topic', '').lower().split()
        keywords.update([w for w in topic_words if len(w) > 3])
        
        # Extract from technical details
        for detail in research_data.get('technicalDetails', []):
            aspect_words = detail.get('aspect', '').lower().split()
            keywords.update([w for w in aspect_words if len(w) > 3])
        
        # Common stop words to remove
        stop_words = {'with', 'using', 'based', 'system', 'approach', 'method'}
        keywords = keywords - stop_words
        
        return list(keywords)[:8]  # Return top 8 keywords
    
    def _generate_paper_content(self, research_data: Dict) -> str:
        """Generate main paper content"""
        sections = []
        
        # Introduction
        sections.append(r"\section{Introduction}")
        sections.append(self._generate_introduction(research_data))
        
        # Related Work
        sections.append(r"\section{Related Work}")
        sections.append(self._generate_related_work(research_data))
        
        # Methodology
        sections.append(r"\section{Methodology}")
        sections.append(self._generate_methodology(research_data))
        
        # Implementation
        sections.append(r"\section{Implementation}")
        sections.append(self._generate_implementation(research_data))
        
        # Evaluation
        sections.append(r"\section{Evaluation}")
        sections.append(self._generate_evaluation(research_data))
        
        # Conclusion
        sections.append(r"\section{Conclusion}")
        sections.append(self._generate_conclusion(research_data))
        
        return '\n\n'.join(sections)
    
    def _generate_introduction(self, research_data: Dict) -> str:
        """Generate introduction section"""
        intro = []
        
        # Background
        intro.append(research_data.get('synthesis', '')[:500] + '...')
        
        # Problem statement
        intro.append(r"\subsection{Problem Statement}")
        if research_data.get('gaps'):
            intro.append("The following research gaps have been identified:")
            intro.append(r"\begin{itemize}")
            for gap in research_data['gaps'][:3]:
                intro.append(f"\\item {gap}")
            intro.append(r"\end{itemize}")
        
        # Contributions
        intro.append(r"\subsection{Contributions}")
        intro.append("This paper makes the following contributions:")
        intro.append(r"\begin{enumerate}")
        intro.append(r"\item A novel approach to " + research_data.get('topic', ''))
        intro.append(r"\item Comprehensive evaluation demonstrating effectiveness")
        intro.append(r"\item Open-source implementation and reproducible results")
        intro.append(r"\end{enumerate}")
        
        return '\n'.join(intro)
    
    def _generate_related_work(self, research_data: Dict) -> str:
        """Generate related work section"""
        related = []
        
        papers = research_data.get('papers', [])[:10]
        
        if papers:
            for paper in papers:
                cite_key = self._generate_cite_key(paper)
                related.append(
                    f"{paper['authors'][0].split()[-1]} et al. \\cite{{{cite_key}}} "
                    f"presented work on {paper['title'][:50]}..."
                )
        
        return '\n\n'.join(related)
    
    def _generate_methodology(self, research_data: Dict) -> str:
        """Generate methodology section"""
        method = []
        
        method.append(research_data.get('proposedApproach', ''))
        
        # Add technical details
        for detail in research_data.get('technicalDetails', []):
            method.append(f"\\subsection{{{detail['aspect']}}}")
            method.append(detail['description'])
            
            if detail.get('implementation'):
                method.append(r"\begin{lstlisting}[language=Python]")
                method.append(detail['implementation'][:200] + "...")
                method.append(r"\end{lstlisting}")
        
        return '\n\n'.join(method)
    
    def _generate_implementation(self, research_data: Dict) -> str:
        """Generate implementation section"""
        impl = []
        
        impl.append("The system was implemented using modern software engineering practices.")
        
        if 'development' in research_data:
            dev = research_data['development']
            impl.append(f"The codebase consists of {dev.get('files', 0)} files "
                       f"with {dev.get('lines', 0)} lines of code.")
            
            impl.append(r"\subsection{Architecture}")
            impl.append("The system architecture follows a modular design pattern...")
        
        return '\n\n'.join(impl)
    
    def _generate_evaluation(self, research_data: Dict) -> str:
        """Generate evaluation section"""
        eval_content = []
        
        eval_content.append(r"\subsection{Experimental Setup}")
        eval_content.append("Experiments were conducted to evaluate the proposed approach...")
        
        if 'testResults' in research_data:
            eval_content.append(r"\subsection{Results}")
            test = research_data['testResults']
            eval_content.append(f"A total of {test.get('total', 0)} tests were executed "
                              f"with {test.get('passed', 0)} passing successfully.")
        
        return '\n\n'.join(eval_content)
    
    def _generate_conclusion(self, research_data: Dict) -> str:
        """Generate conclusion section"""
        return "This paper presented a novel approach to " + research_data.get('topic', '') + \
               ". Future work includes extending the approach to additional domains."
    
    def _generate_bibliography(self, paper_dir: Path, papers: List[Dict]):
        """Generate BibTeX bibliography"""
        bib_content = []
        
        for paper in papers[:30]:  # Limit to 30 references
            cite_key = self._generate_cite_key(paper)
            
            if paper['source'] == 'arxiv':
                entry = f"""@article{{{cite_key},
  title={{{paper['title']}}},
  author={{{' and '.join(paper['authors'][:5])}}},
  journal={{arXiv preprint arXiv:{paper.get('arxiv_id', '')}}},
  year={{{paper.get('year', 2024)}}}
}}"""
            else:
                entry = f"""@inproceedings{{{cite_key},
  title={{{paper['title']}}},
  author={{{' and '.join(paper['authors'][:5])}}},
  booktitle={{{paper.get('venue', 'Proceedings')}}},
  year={{{paper.get('year', 2024)}}}
}}"""
            
            bib_content.append(entry)
        
        bib_file = paper_dir / 'references.bib'
        with open(bib_file, 'w') as f:
            f.write('\n\n'.join(bib_content))
    
    def _generate_cite_key(self, paper: Dict) -> str:
        """Generate citation key for paper"""
        first_author = paper['authors'][0].split()[-1].lower() if paper['authors'] else 'unknown'
        year = paper.get('year', 2024)
        title_word = paper['title'].split()[0].lower()
        return f"{first_author}{year}{title_word}"
    
    async def generate_presentation(
        self,
        session_id: str,
        research_data: Dict,
        style: str = 'modern'
    ) -> str:
        """Generate PowerPoint presentation"""
        try:
            # Create presentation
            prs = Presentation()
            
            # Define slide layouts
            title_slide_layout = prs.slide_layouts[0]
            content_slide_layout = prs.slide_layouts[1]
            section_slide_layout = prs.slide_layouts[2]
            
            # Title slide
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title.text = research_data.get('topic', 'Research Presentation')
            subtitle.text = f"Generated on {datetime.now().strftime('%B %d, %Y')}"
            
            # Agenda slide
            slide = prs.slides.add_slide(content_slide_layout)
            slide.shapes.title.text = "Agenda"
            content = slide.placeholders[1]
            agenda_items = [
                "Introduction & Background",
                "Research Objectives",
                "Literature Review",
                "Methodology",
                "Implementation",
                "Results & Evaluation",
                "Conclusion & Future Work"
            ]
            content.text = '\n'.join([f"• {item}" for item in agenda_items])
            
            # Introduction slides
            self._add_introduction_slides(prs, research_data)
            
            # Research objectives
            self._add_objectives_slide(prs, research_data)
            
            # Literature review
            self._add_literature_slides(prs, research_data)
            
            # Methodology
            self._add_methodology_slides(prs, research_data)
            
            # Implementation
            self._add_implementation_slides(prs, research_data)
            
            # Results
            self._add_results_slides(prs, research_data)
            
            # Conclusion
            self._add_conclusion_slide(prs, research_data)
            
            # Thank you slide
            slide = prs.slides.add_slide(title_slide_layout)
            slide.shapes.title.text = "Thank You"
            slide.placeholders[1].text = "Questions?"
            
            # Apply styling
            self._apply_presentation_style(prs, style)
            
            # Save presentation
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            filename = f"presentation_{session_id}_{timestamp}.pptx"
            filepath = OUTPUT_DIR / filename
            prs.save(str(filepath))
            
            return str(filepath)
            
        except Exception as e:
            raise Exception(f"Error generating presentation: {str(e)}")
    
    def _add_introduction_slides(self, prs: Presentation, research_data: Dict):
        """Add introduction slides"""
        # Background slide
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Background"
        
        content = slide.placeholders[1]
        background_text = research_data.get('synthesis', '')[:300] + '...'
        content.text = background_text
        
        # Problem statement
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Problem Statement"
        
        content = slide.placeholders[1]
        gaps = research_data.get('gaps', [])
        if gaps:
            content.text = "Key Challenges:\n" + '\n'.join([f"• {gap}" for gap in gaps[:4]])
    
    def _add_objectives_slide(self, prs: Presentation, research_data: Dict):
        """Add research objectives slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Research Objectives"
        
        content = slide.placeholders[1]
        objectives = research_data.get('objectives', [])
        if objectives:
            content.text = '\n'.join([f"• {obj}" for obj in objectives[:5]])
    
    def _add_literature_slides(self, prs: Presentation, research_data: Dict):
        """Add literature review slides"""
        # Overview slide
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Literature Review"
        
        content = slide.placeholders[1]
        papers = research_data.get('papers', [])
        
        if papers:
            review_text = f"Analyzed {len(papers)} relevant papers\n\n"
            review_text += "Key contributions from literature:\n"
            
            for paper in papers[:3]:
                review_text += f"• {paper['authors'][0].split()[-1]} et al. ({paper.get('year', 'n.d.')}): "
                review_text += f"{paper['title'][:50]}...\n"
            
            content.text = review_text
    
    def _add_methodology_slides(self, prs: Presentation, research_data: Dict):
        """Add methodology slides"""
        # Overview
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Methodology Overview"
        
        content = slide.placeholders[1]
        approach = research_data.get('proposedApproach', '')
        content.text = approach[:400] + '...' if len(approach) > 400 else approach
        
        # Technical details
        for detail in research_data.get('technicalDetails', [])[:3]:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = detail['aspect']
            
            content = slide.placeholders[1]
            content.text = detail['description'][:300] + '...'
    
    def _add_implementation_slides(self, prs: Presentation, research_data: Dict):
        """Add implementation slides"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Implementation"
        
        content = slide.placeholders[1]
        
        if 'development' in research_data:
            dev = research_data['development']
            impl_text = f"Project Statistics:\n"
            impl_text += f"• Files: {dev.get('files', 0)}\n"
            impl_text += f"• Lines of Code: {dev.get('lines', 0)}\n"
            impl_text += f"• Test Coverage: {dev.get('testCoverage', 0)}%\n"
            impl_text += f"• Components: {dev.get('components', 0)}"
            
            content.text = impl_text
    
    def _add_results_slides(self, prs: Presentation, research_data: Dict):
        """Add results slides"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Results & Evaluation"
        
        content = slide.placeholders[1]
        
        if 'testResults' in research_data:
            test = research_data['testResults']
            results_text = f"Testing Results:\n"
            results_text += f"• Total Tests: {test.get('total', 0)}\n"
            results_text += f"• Passed: {test.get('passed', 0)}\n"
            results_text += f"• Failed: {test.get('failed', 0)}\n"
            results_text += f"• Coverage: {test.get('coverage', 0)}%"
            
            content.text = results_text
    
    def _add_conclusion_slide(self, prs: Presentation, research_data: Dict):
        """Add conclusion slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Conclusion & Future Work"
        
        content = slide.placeholders[1]
        
        conclusion_text = "Key Contributions:\n"
        conclusion_text += f"• Novel approach to {research_data.get('topic', '')}\n"
        conclusion_text += "• Comprehensive implementation and evaluation\n"
        conclusion_text += "• Open-source release\n\n"
        conclusion_text += "Future Directions:\n"
        conclusion_text += "• Extend to additional domains\n"
        conclusion_text += "• Performance optimization\n"
        conclusion_text += "• Real-world deployment"
        
        content.text = conclusion_text
    
    def _apply_presentation_style(self, prs: Presentation, style: str):
        """Apply styling to presentation"""
        # Modern style with dark theme
        if style == 'modern':
            for slide in prs.slides:
                # Set background
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(30, 30, 30)  # Dark background
                
                # Style title
                if slide.shapes.title:
                    title = slide.shapes.title
                    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                    title.text_frame.paragraphs[0].font.size = Pt(32)
                    title.text_frame.paragraphs[0].font.bold = True
                
                # Style content
                for shape in slide.shapes:
                    if shape.has_text_frame and shape != slide.shapes.title:
                        for paragraph in shape.text_frame.paragraphs:
                            paragraph.font.color.rgb = RGBColor(220, 220, 220)
                            paragraph.font.size = Pt(18)

# Flask routes
generator = DocumentGenerator()

@app.route('/health', methods=['GET'])
def health_check():
    return jsonify({'status': 'healthy'})

@app.route('/generate/report', methods=['POST'])
async def generate_report():
    """Generate PDF report"""
    try:
        data = request.json
        session_id = data['sessionId']
        
        # Get research data from Redis
        research_data = redis_client.get(f'research:results:{session_id}')
        if not research_data:
            return jsonify({'error': 'Research data not found'}), 404
        
        research_data = json.loads(research_data)
        
        # Generate report
        filepath = await generator.generate_pdf_report(session_id, research_data)
        
        return jsonify({
            'success': True,
            'filepath': filepath,
            'filename': os.path.basename(filepath)
        })
        
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/generate/paper', methods=['POST'])
async def generate_paper():
    """Generate LaTeX paper"""
    try:
        data = request.json
        session_id = data['sessionId']
        template = data.get('template', 'ieee')
        
        # Get research data
        research_data = redis_client.get(f'research:results:{session_id}')
        if not research_data:
            return jsonify({'error': 'Research data not found'}), 404
        
        research_data = json.loads(research_data)
        
        # Generate paper
        filepath = await generator.generate_latex_paper(session_id, research_data, template)
        
        return jsonify({
            'success': True,
            'filepath': filepath,
            'filename': os.path.basename(filepath)
        })
        
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/generate/presentation', methods=['POST'])
async def generate_presentation():
    """Generate PowerPoint presentation"""
    try:
        data = request.json
        session_id = data['sessionId']
        style = data.get('style', 'modern')
        
        # Get research data
        research_data = redis_client.get(f'research:results:{session_id}')
        if not research_data:
            return jsonify({'error': 'Research data not found'}), 404
        
        research_data = json.loads(research_data)
        
        # Generate presentation
        filepath = await generator.generate_presentation(session_id, research_data, style)
        
        return jsonify({
            'success': True,
            'filepath': filepath,
            'filename': os.path.basename(filepath)
        })
        
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/download/<filename>', methods=['GET'])
def download_file(filename):
    """Download generated file"""
    filepath = OUTPUT_DIR / filename
    if filepath.exists():
        return send_file(str(filepath), as_attachment=True)
    else:
        return jsonify({'error': 'File not found'}), 404

@app.route('/templates', methods=['GET'])
def list_templates():
    """List available templates"""
    templates = [f.stem for f in TEMPLATES_DIR.glob('*_template.tex')]
    return jsonify({
        'templates': templates,
        'default': 'ieee'
    })

if __name__ == '__main__':
    port = int(os.environ.get('PORT', 5001))
    app.run(host='0.0.0.0', port=port, debug=False)